import os

from .base import BaseParser
from .core import has_cmd, logger, run_command


class PPTXParser(BaseParser):
    def parse(self):
        logger.info(f"[PPTX] Processing: {self.basename}")

        text_extracted = False
        from .core import get_apryse_args

        # 1. Apryse docpub (Highest quality conversion)
        if has_cmd("docpub"):
            logger.info("    → docpub (Apryse)")
            args = get_apryse_args("docpub")
            res = run_command(args + ["-f", "pdf", "-o", self.norm_dir, self.file_path])
            pdf_out = os.path.join(
                self.norm_dir, os.path.splitext(self.basename)[0] + ".pdf"
            )
            if res and res.returncode == 0 and os.path.exists(pdf_out):
                if has_cmd("pdf2text"):
                    args_p2t = get_apryse_args("pdf2text")
                    run_command(
                        args_p2t + ["-f", "plain", "-o", self.norm_dir, pdf_out]
                    )
                    apryse_txt = os.path.join(
                        self.norm_dir, os.path.splitext(self.basename)[0] + ".txt"
                    )
                    if os.path.exists(apryse_txt):
                        os.rename(apryse_txt, self.content_txt_path)
                        text_extracted = True
                elif has_cmd("pdftotext"):
                    run_command(
                        ["pdftotext", "-layout", pdf_out, self.content_txt_path]
                    )
                    text_extracted = True

        # 2. python-pptx
        if not text_extracted:
            try:
                from pptx import Presentation

                prs = Presentation(self.file_path)
                out = []
                for i, slide in enumerate(prs.slides, 1):
                    out.append(f"--- Slide {i} ---")
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text:
                            out.append(shape.text)
                    out.append("")
                self.write_text("\n".join(out))
                text_extracted = True
            except ImportError:
                logger.info("    (!) python-pptx not installed")
            except Exception as e:
                logger.warning(f"python-pptx error: {e}")

        # 3. LibreOffice Fallback
        if not text_extracted and (has_cmd("soffice") or has_cmd("libreoffice")):
            soffice = "soffice" if has_cmd("soffice") else "libreoffice"
            run_command(
                [
                    soffice,
                    "--headless",
                    "--convert-to",
                    "pdf",
                    "--outdir",
                    self.norm_dir,
                    self.file_path,
                ]
            )
            pdf_file = os.path.join(
                self.norm_dir, os.path.splitext(self.basename)[0] + ".pdf"
            )
            if os.path.exists(pdf_file) and has_cmd("pdftotext"):
                run_command(["pdftotext", "-layout", pdf_file, self.content_txt_path])
                text_extracted = True
                if has_cmd("pdftoppm"):
                    run_command(
                        [
                            "pdftoppm",
                            "-png",
                            "-r",
                            "180",
                            pdf_file,
                            os.path.join(self.preview_dir, "slide"),
                        ]
                    )
                os.remove(pdf_file)

        if not text_extracted:
            self.write_text(f"[No text extracted for {self.basename}]")

        self.generate_markdown_wrapper("Presentation")
        logger.info(f"✓ PPTX parsed: {self.basename}")
